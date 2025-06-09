from pptx import Presentation
from pptx.util import Inches

# 1. 새 프레젠테이션 또는 기존 템플릿 불러오기
prs = Presentation('your_template.pptx')  # 템플릿이 없다면 Presentation() 사용

# 2. 이미지 파일 리스트 준비
image_files = [
    'pairplot_result1.png',
    'pairplot_result2.png',
    'pairplot_result3.png'
]

# 3. 반복문으로 각 이미지를 새 슬라이드에 삽입
for img_path in image_files:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(5)
    slide.shapes.add_picture(img_path, left, top, width, height)

# 4. 결과 저장
prs.save('결과_시각화_일괄삽입.pptx')
